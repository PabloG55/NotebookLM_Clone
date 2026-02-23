"""
Handles loading text from PDF, PPTX, TXT files and URLs.
Returns raw text string.
"""
import io
import fitz  # PyMuPDF
from pptx import Presentation
import requests
from bs4 import BeautifulSoup


def load_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text.strip()


def load_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
    return text.strip()


def load_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8").strip()
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1").strip()


def load_url(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0"}
    response = requests.get(url, headers=headers, timeout=15)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    # Remove scripts and styles
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Clean up blank lines
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


def ingest_source(source_type: str, data) -> str:
    """
    source_type: 'pdf', 'pptx', 'txt', 'url'
    data: bytes for files, str for url
    """
    if source_type == "pdf":
        return load_pdf(data)
    elif source_type == "pptx":
        return load_pptx(data)
    elif source_type == "txt":
        return load_txt(data)
    elif source_type == "url":
        return load_url(data)
    else:
        raise ValueError(f"Unsupported source type: {source_type}")
